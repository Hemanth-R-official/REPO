import pandas as pd
import matplotlib.pyplot as plt
import os
from tkinter import filedialog
import tkinter as tk
import tkinter.messagebox
from pptx import Presentation  
from pptx.util import Inches,Pt  
import smtplib 
from email.mime.multipart import MIMEMultipart 
from email.mime.text import MIMEText 
from email.mime.base import MIMEBase 
from email import encoders 

#***************************************************************************************************************************************************************************************************************************
def locate():
	filepath = filedialog.askopenfilename()
	if filepath.endswith('.xlsx'):
		return filepath 
	else:
		tkinter.messagebox.showinfo("Alert !", "Unsupported File...")
		response=tkinter.messagebox.askquestion("Unsupported File","Do you want to relocate it ? or return to main menu.")
		if response=='yes':
			locate()	
		else:
			exit()
#*****************************************************************************************************************************************************************************************************************
def ppt(REG,N,L,Q,V,G1,G2,M,SA,CGPA,RC,SL):

    address=""    
    address+=os.getcwd()
    logo= address+'\\'+'frontpage.jpg'
    reffileaddress=address+'\\'+"REFERENCE.txt"
    reffile= open(reffileaddress,"r")
    ref=""
    for i in reffile:
        ref+=i 
	
    ppt = Presentation()
	
#@@@@@@@@@@@@@@@@@@@@@@@ PAGE 1 @@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@

    blank_slide_layout = ppt.slide_layouts[6]  
    slide = ppt.slides.add_slide(blank_slide_layout)  
	
    left =  Inches(1.5)
    top = Inches(0)
	
    pic = slide.shapes.add_picture(logo, 
								left, top) 
    left = Inches(1)  
    height = Inches(0)
    pic = slide.shapes.add_picture(logo, left, 
								top, height = height)	

	#TEXT edit  
    left = Inches(3) 
    top = Inches(4.5)
    width = height = Inches(4)  

	# creating textBox 
    txBox = slide.shapes.add_textbox(left, top, 
									width, height) 
	
	# creating textFrames 
    tf = txBox.text_frame 
    p = tf.add_paragraph()
    p.text = "NAME : "+N
    p.font.bold = True
    p.font.size = Pt(25)

    p = tf.add_paragraph()
    p.text = "REG NO : "+str(REG)
    p.font.bold = True
    p.font.size = Pt(25)

#@@@@@@@@@@@@@@@@@@@@@@@@@@ PAGE 2 @@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@

    graph1= address+'\\'+G1
	# Selecting blank slide 
    blank_slide_layout = ppt.slide_layouts[6]  
	
	# Attaching slide to ppt 
    slide = ppt.slides.add_slide(blank_slide_layout)  
	
	# For margins 
    left = Inches(0)
    top = Inches(0)
	
	# adding images 
    pic = slide.shapes.add_picture(graph1, 
								left, top) 
	
    left = Inches(1)  
    height = Inches(0) #1 	
    pic = slide.shapes.add_picture(graph1, left, 
								top, height = height)

	#text
    left = Inches(16.5/2.54)
    top = Inches(2.5/2.54)
    height = Inches(5.7/2.54)  
    width =Inches(5.2/2.54)   
    txBox = slide.shapes.add_textbox(left, top, 
									width, height) 
    tf = txBox.text_frame 
    p = tf.add_paragraph()
    p.text ="QUANTS : "+str(Q)
    p.font.bold = True
    p.font.size = Pt(25)

    p = tf.add_paragraph()
    p.text ="   "

    p = tf.add_paragraph()
    p.text ="LOGICAL : "+str(L)
    p.font.bold = True
    p.font.size = Pt(25)

    p = tf.add_paragraph()
    p.text ="   "	

    p = tf.add_paragraph()
    p.text ="VERBAL  : "+str(V)
    p.font.bold = True
    p.font.size = Pt(25)

    left = Inches(2/2.54) 
    top = Inches(12/2.54)
    height = Inches(2.3/2.54)   
    width = Inches(19/2.54)
    txBox = slide.shapes.add_textbox(left, top, 
									width, height) 
    tf = txBox.text_frame 
    p = tf.add_paragraph()
    avg=Q+L+V
    if(avg>=60):
        p.text="* EXCELLENT ! Keep it up..."
    elif(avg<=30):
        p.text="* AVERAGE ! Danger Zone, Need more practice..."
    else:
        p.text="* GOOD ! Take a few more tests to improve..."
    p.font.size = Pt(25)

    left = Inches(2/2.54) 
    top = Inches(13.5/2.54)
    width = Inches(19/2.54)
    height = Inches(2.3/2.54)   
    txBox = slide.shapes.add_textbox(left, top, 
									width, height) 
    tf = txBox.text_frame 
    p = tf.add_paragraph()
    second_text="* Try to improve more on "
    f=0
    if(Q<=7 or L<=7 or V<=7):
        if(Q<=7):
            second_text+=' quants'
            f+=1
        if(L<=7):
            if(f!=0):
                second_text+=','
            second_text+=' logical'
            f+=1
        if(V<=7):
            if(f>=1):
                second_text+=' and'
            second_text+=' verbal'
        second_text+='.'
        p.text=second_text
        p.font.size = Pt(25)
    else:
        if(Q<L and Q<V):
            second_text+=' quants.'
        elif(L<Q and L<V):
            second_text+=' logical.'
        else:
            second_text+=' verbal.'
        p.text=second_text
        p.font.size = Pt(25)
    

#@@@@@@@@@@@@@@@@@@@@@@@@@@ PAGE 3 @@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@
#SECOND GRAPH
	
    graph2= address+'\\'+G2
    blank_slide_layout = ppt.slide_layouts[6]  
	
	# Attaching slide to ppt 
    slide = ppt.slides.add_slide(blank_slide_layout)  
	
	# For margins 
    left = Inches(0)
    top = Inches(0)
	
	# adding images 
    pic = slide.shapes.add_picture(graph2, 
								left, top) 
    left = Inches(0)  
    height = Inches(0)	
    pic = slide.shapes.add_picture(graph2, left, 
								top, height = height)


	#text
    left = Inches(15/2.54) 
    top = Inches(2.5/2.54)
    height = Inches(6.6/2.54) 
    width =Inches(7.6/2.54)   
    txBox = slide.shapes.add_textbox(left, top, 
									width, height) 
    tf = txBox.text_frame 

    p = tf.add_paragraph()
    p.text ="SKILLRACK LEVEL : "+str(SL)
    p.font.bold = True
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text ="   "

    p = tf.add_paragraph()
    p.text ="RESUME COUNT : "+str(RC)
    p.font.bold = True
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text ="   "

    p = tf.add_paragraph()
    p.text ="CGPA : "+str(CGPA)
    p.font.bold = True
    p.font.size = Pt(20)

    p = tf.add_paragraph()
    p.text ="   "

    p = tf.add_paragraph()
    p.text ="STANDING ARREARS : "+str(SA)
    p.font.bold = True
    p.font.size = Pt(20)

	#placement eligiblity 
    left = Inches(2/2.54) 
    top = Inches(12/2.54)
    height = Inches(2.3/2.54)
    width = Inches(19/2.54) 
    txBox = slide.shapes.add_textbox(left, top, 
									width, height) 
    tf = txBox.text_frame 
    if(SA<=1 and CGPA>6):
        p = tf.add_paragraph()
        p.text ="* Eligible for placements ! CONGRATULATIONS..."
        p.font.size = Pt(25)

    else:
        p = tf.add_paragraph()
        p.text ="* Not eligible for placements ! NEED TO WORK HARD..."
        p.font.size = Pt(25)
     

#@@@@@@@@@@@@@@@@@@@@@@@@@@ PAGE 4 @@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@
#REFFERENCE

    blank_slide_layout = ppt.slide_layouts[6]
    slide = ppt.slides.add_slide(blank_slide_layout)

    left = Inches(2.5) 
    top = Inches(0)
    width = height = Inches(10)  

    txBox = slide.shapes.add_textbox(left, top, 
									width, height) 
	
	# creating textFrames 
    tf = txBox.text_frame 
    p = tf.add_paragraph()
    p.text ="REFERENCE LINKS"
    p.font.bold = True
    p.font.size = Pt(40)

	#linktext

    left = Inches(0.5) 
    top = Inches(1)
    width = height = Inches(10)  
    txBox = slide.shapes.add_textbox(left, top, 
									width, height) 
    tf = txBox.text_frame 
    p = tf.add_paragraph()
    p.text =ref
    p.font.size = Pt(15)
	
#@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@@
    Name=N
    Fname=Name+"'s Report"+".pptx"
    ppt.save(Fname)
    pptaddress=mail(Name,M)
    return(pptaddress)

#**************************************************************************************************************************************************************************************************************
def mail(N,M):
	add=""
	add+=os.getcwd()
	add+="\\"
	ext=N+"'s Report"+".pptx"
	address=add+ext

	fromaddr ='' # ENTER MAIL ID!!!!!!!!!!!!!!
	toaddr = M

	# instance of MIMEMultipart 
	msg = MIMEMultipart() 
	
	# senders email address   
	msg['From'] = fromaddr 
	
	# receivers email address  
	msg['To'] = toaddr 
	
	# storing the subject  
	msg['Subject'] = "REPO"
	
	# string to store the body of the mail 
	body = "The greatest glory in living lies not in never falling, but in rising every time we fall. -Nelson Mandela"
	
	# body with the msg instance 
	msg.attach(MIMEText(body, 'plain')) 
	
	# open the file to be sent  
	filename = "Repo.pptx"
	attachment = open(address, "rb") 

	# instance of MIMEBase and named as p 
	p = MIMEBase('application', 'octet-stream') 
	
	# To change the payload into encoded form 
	p.set_payload((attachment).read()) 
	
	# encode into base64 
	encoders.encode_base64(p) 
	
	p.add_header('Content-Disposition', "attachment; filename= %s" % filename) 
	
	# attach the instance 'p' to instance 'msg' 
	msg.attach(p) 
	
	# creates SMTP session 
	s = smtplib.SMTP('smtp.gmail.com', 587) 
	
	# start TLS for security 
	s.starttls() 
	
	# Authentication 
	s.login(fromaddr, "") #ENTER MAIL PASSWORD !!!!!    WARNING...................
	
	# Converts the Multipart msg into a string 
	text = msg.as_string() 
	
	# sending the mail 
	s.sendmail(fromaddr, toaddr, text) 
	
	# terminating the session 
	s.quit() 

	#delete files
	return(address)

#**************************************************************************************************************************************************************************************************************************
def run():
    file_path=locate()
    data=pd.read_excel(
        os.path.join(file_path),
        engine='openpyxl',
    )
    a=data['S.NO']
    
    for i in range(-1,len(a)-1):
        p=i+1
        Q=data.at[p,'QUANTS']
        M=data.at[p,'EMAIL']
        L=data.at[p,'LOGICAL']
        V=data.at[p,'VERBAL']
        N=data.at[p,'NAME']
        C=data.at[p,'CORRECT']
        W=data.at[p,'WRONG']
        SA=data.at[p,'STANDING ARREARS']
        RC=data.at[p,'RESUME COUNT']
        CGPA=data.at[p,'CGPA']
        SL=	data.at[p,'SKILLRACK LEVEL']
        Graphname1=N+"'s 1 graph"+'.png'
        Graphname2=N+"'s 2 graph"+'.png'
        x=['QUANTS','LOGICAL','VERBAL']
        y=[Q,L,V]

        f1=plt.figure(1)
        plt.bar(x,y,color=['red','blue','green'])
        plt.title(N+"'s Technical Skills Report")
        plt.savefig(Graphname1)
        plt.clf()

        f2=plt.figure(2)
        y=[C,W]
        pielables=["Correct","Wrong"]
        pieexplode=[0.2,0]
        plt.pie(y,labels=pielables,explode=pieexplode,shadow=True)
        plt.title("Correct VS Wrong Answers")
        plt.savefig(Graphname2)
        plt.clf()
        plt.close('all')
			
        pptaddress=ppt(p+1,N,L,Q,V,Graphname1,Graphname2,M,SA,CGPA,RC,SL)
        os.remove(Graphname1)
        os.remove(Graphname2)
        os.remove(pptaddress)

    tkinter.messagebox.showinfo("Done !", "Task accomplished....")
        

def main():
	run()

if __name__=="__main__":
    main()
